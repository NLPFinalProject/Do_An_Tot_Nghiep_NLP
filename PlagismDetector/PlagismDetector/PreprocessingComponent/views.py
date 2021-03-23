import tika
from tika import parser
import re
import os
import sys
import win32com.client as win32
from win32com.client import constants
from pptx import Presentation
import pandas as pd
import openpyxl
from vncorenlp import VnCoreNLP
import zipfile
import xml.etree.ElementTree
import tika
import re
from tika import parser
import logging
from docx_utils.flatten import opc_to_flat_opc
from xml.dom import minidom
import os.path
import win32com.client
import docx2txt
import pythoncom
# LƯU Ý:
# input của các hàm xử lý là tên file và output của các hàm xử lý ...2txt là một list các câu.
# từng phần tử của 1 list lớn (1 câu) là 1 list nhỏ chứa các từ được phân tách nhờ VNCORE.

#### ------------------------------------các function hỗ trợ cho function docx2txt
# Get paragraph string. Input is paragraph element
def para_string(para):
    string = ""
    # if (str(para)[21:34] not in str(wp_tbl)):# and (str(para)[21:34] not in str(wp_txbx)):
    wt = para.getElementsByTagName('w:t')
    for i in range(len(wt)):
        string = string + wt[i].firstChild.data

    return string


# Get table string. Input is table element
def table_string(table):
    string = ""

    wp = table.getElementsByTagName('w:p')
    column = len(table.getElementsByTagName('w:tc')) / len(table.getElementsByTagName('w:tr'))
    c = 1
    for i in range(len(wp)):
        string = string + para_string(wp[i])
        if c % column == 0 and c != len(wp):
            string = string + '\n'
        else:
            string = string + '\t'
        c = c + 1
    return string


# Get all elements
def get_all_elements(lst, type_of_element):
    elements_list = []
    for i in range(len(lst)):
        Elements = lst[i].getElementsByTagName(type_of_element)
        for elm in Elements:
            elements_list.append(elm)

    return elements_list


##--------------------------------------------------------------------

def docx2text(docx_file_name):
    text=docx2txt.process(docx_file_name)
    print("path------------------",sys.path[-1])
    #vncorenlp_file = r'VnCoreNLP\VnCoreNLP-1.1.1.jar'
    #path=sys.path[-1]+"\\PreprocessingComponent\\VnCoreNLP\\VnCoreNLP-1.1.1.jar"
    path= "D:\\study\\PlagismDetector\\PlagismDetector\\PreprocessingComponent\\VnCoreNLP\\VnCoreNLP-1.1.1.jar"
    vncorenlp_file= path
    print("path:                        ", path)
    with VnCoreNLP(vncorenlp_file, annotators="wseg", max_heap_size='-Xmx4g', quiet=False) as vncorenlp:
        split_sentence = vncorenlp.tokenize(text)
    return split_sentence

# Rút trích câu từ file docx.
# def docx2txt(docx_file_name):
#     # Parse xml file
#     xml_file_name = 'mydocx.xml'
#     opc_to_flat_opc(docx_file_name, xml_file_name)
#     my_docx = minidom.parse(xml_file_name)

#     # Get elements
#     paragraph = my_docx.getElementsByTagName('w:p')
#     table = my_docx.getElementsByTagName('w:tbl')

#     # Get all w:p elements in table elements. Output is two-dimensional list
#     wp_tbl = get_all_elements(table, 'w:p')

#     # Get text and save to "string" variable
#     para_index = 0
#     tbl_index = 0
#     string = ""
#     while para_index < len(paragraph):
#         if paragraph[para_index] in wp_tbl:
#             string = string + table_string(table[tbl_index])
#             para_index += len(table[tbl_index].getElementsByTagName('w:p'))
#             tbl_index += 1
#         else:
#             string = string + para_string(paragraph[para_index])
#             para_index += 1
#         # string = string + '.'

#     #vncorenlp_file = r'VnCoreNLP\VnCoreNLP-1.1.1.jar'
#     path=sys.path[-1]+"\\VnCoreNLP\\VnCoreNLP-1.1.1.jar"
#     vncorenlp_file= path
#     print("path:                        ", path)
#     # with VnCoreNLP(vncorenlp_file, annotators="wseg", max_heap_size='-Xmx4g', quiet=False) as vncorenlp:
#     #     split_sentence = vncorenlp.tokenize(string)
#     return path


# hàm chuyển đổi định dạng từ fild .doc sang .docx
def doc2docx(filename, path=os.getcwd()):
    # Dùng khi chuyển tất cả các file doc sang docx có trong đường dẫn baseDir

    # baseDir = os.path.abspath(os.getcwd())  # Starting directory for directory walk
    # word = win32com.client.Dispatch("Word.application")
    # for dir_path, dirs, files in os.walk(baseDir):
    #     for file_name in files:
    #
    #         file_path = os.path.join(dir_path, file_name)
    #         file_name, file_extension = os.path.splitext(file_path)
    #
    #         if "~$" not in file_name:
    #             if file_extension.lower() == '.doc':  #
    #                 docx_file = '{0}{1}'.format(file_path, 'x')
    #
    #                 if not os.path.isfile(docx_file):  # Skip conversion where docx file already exists
    #
    #                     file_path = os.path.abspath(file_path)
    #                     docx_file = os.path.abspath(docx_file)
    #                     try:
    #                         wordDoc = word.Documents.Open(file_path)
    #                         wordDoc.SaveAs2(docx_file, FileFormat=16)
    #                         wordDoc.Close()
    #                     except Exception as e:
    #                         print('Failed to Convert: {0}'.format(file_path))
    #                         print(e)

    baseDir = os.path.abspath(os.getcwd())  # Starting directory for directory walk
    ###testing mode
    pythoncom.CoInitialize()
    word = win32com.client.Dispatch("Word.application")
    pythoncom.CoInitialize()
    ###end test
    file_path = os.path.join(baseDir, filename)
    file_name, file_extension = os.path.splitext(file_path)

    if "~$" not in file_name:
        if file_extension.lower() == '.doc':  #
            docx_file = '{0}{1}'.format(file_path, 'x')

            if not os.path.isfile(docx_file):  # Skip conversion where docx file already exists

                file_path = os.path.abspath(file_path)
                docx_file = os.path.abspath(docx_file)
                try:
                    wordDoc = word.Documents.Open(file_path)
                    wordDoc.SaveAs2(docx_file, FileFormat=16)
                    wordDoc.Close()
                except Exception as e:
                    print('Failed to Convert: {0}'.format(file_path))
                    print(e)


# hàm rút trích các câu từ file ppt
def ppt2txt(filename):
    ppt = Presentation(filename)
    sentences = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                sentences += shape.text + ". "

    #vncorenlp_file = r'VnCoreNLP\VnCoreNLP-1.1.1.jar'
    path=sys.path[-1]+"\\VnCoreNLP\\VnCoreNLP-1.1.1.jar"
    vncorenlp_file= path
    
    with VnCoreNLP(vncorenlp_file, annotators="wseg", max_heap_size='-Xmx4g', quiet=False) as vncorenlp:
        split_sentence = vncorenlp.tokenize(sentences)
    return split_sentence


# hàm rút trich các câu từ file excel ### cần cập nhật thêm vì chưa hoàn thiện
def xlsx2txt(filename):
    data = pd.read_excel(filename, sheet_name=None, index_col=0, keep_default_na=0)
    list_sheet = []
    for key, value in data.items():
        list_sheet.append(value)
    listToStr = '.'.join(map(str, list_sheet))
    list_sen = listToStr.split('\n')
    for i in range(0, len(list_sen)):
        list_sen[i] = " ".join(list_sen[i].split())
    return list_sen


# hàm rút trích và tách câu từ file pdf.
def pdf2txt(filename):
    parsed = parser.from_file(filename)
    data = parsed["content"]
    list_sen = data.split('\s{4,}')
    # for i in range(0, len(list_sen)):
    sentences = " ".join(list_sen[0].split())
    #vncorenlp_file = r'VnCoreNLP\VnCoreNLP-1.1.1.jar'
    path=sys.path[-1]+"\\VnCoreNLP\\VnCoreNLP-1.1.1.jar"
    vncorenlp_file= path

    with VnCoreNLP(vncorenlp_file, annotators="wseg", max_heap_size='-Xmx4g', quiet=False) as vncorenlp:
        split_sentence = vncorenlp.tokenize(sentences)
    return split_sentences


def convert2listsen(vncoretoken):
    # remove .
    for i in vncoretoken:
        for j in i:
            if(j == "."):
                i.remove(j)
    s = ""
    res=[]
    h=0
    for i in vncoretoken:
        for j in i:
           s += j + " "
        s = s.strip()
        s+="."
        res.append(s)
        s=""
    return res

def preprocess(filename):
    # thực hiện đưa filename nào muốn xử lý vào biến filename bên dưới và đợi kết quả trên màn hình.
    name, file_extension = os.path.splitext(filename)
    print(name)
    print(file_extension)
    if (file_extension.lower() == ".doc"):
        doc2docx(filename)
        a = docx2text(filename + "x")
        b = convert2listsen(a) # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word=[] # số từ của câu đầu tiên tương tự cho a[1],....
        for i in a:
            num_word.append(len(i))

    elif (file_extension.lower() == ".docx"):
        a = docx2text(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = []  # số từ của câu đầu tiên tương tự cho a[1],....
        for i in a:
            num_word.append(len(i))

    elif (file_extension.lower() == ".pdf"):
        a = pdf2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = []  # số từ của câu đầu tiên tương tự cho a[1],....
        for i in a:
            num_word.append(len(i))

    elif (file_extension.lower() == ".xlsx"):
        a = xlsx2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = []  # số từ của câu đầu tiên tương tự cho a[1],....
        for i in a:
            num_word.append(len(i))
    elif (file_extension.lower() == ".pptx"):
        a = ppt2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = []  # số từ của câu đầu tiên tương tự cho a[1],....
        for i in a:
            num_word.append(len(i))
    return filename, b, num_word # filename, list câu. số từ của mỗi câu

def main():
    filename = "D:/kamen rider.doc"
    a = doc2docx(filename)
    print(a)
    b = preprocess(filename)
    print(b)
if __name__=='__main__':
    a = docx2txt("../../ABC/vanbandemo.docx")
    print(a)
#print(b)
# a: filename
# b: list câu
# c: số từ của mỗi câu trong b