import csv
import math
import os
import sys
import time
import uuid
from string import punctuation
from xml.dom import minidom

import pandas as pd
import pythoncom
import win32com.client
from docx_utils.flatten import opc_to_flat_opc
from pptx import Presentation
from underthesea import sent_tokenize, word_tokenize
sys.path.insert(1, os.getcwd() + "/PreprocessingComponent")
from PreprocessingComponent.pdfminer3 import Pdf_extract


# -----các function hỗ trợ cho function DocxToText-
# Get paragraph string. Input is paragraph element
def ParaString(para):
    string = ""
    # if (str(para)[21:34] not in str(wpTbl)):
    # and (str(para)[21:34] not in str(wp_txbx)):
    wt = para.getElementsByTagName('w:t')
    for i in range(len(wt)):
        string = string + wt[i].firstChild.data
    return string


# Get table string. Input is table element
def TableString(table):
    string = ""
    wp = table.getElementsByTagName('w:p')
    temp1 = len(table.getElementsByTagName('w:tc'))
    temp2 = len(table.getElementsByTagName('w:tr'))
    column = temp1 / temp2
    c = 1
    for i in range(len(wp)):
        string = string + ParaString(wp[i])
        if c % column == 0 and c != len(wp):
            string += '. '
        else:
            string += '. '
        c += 1
    return string


# Get all elements
def GetAllElements(lst, type_of_element):
    elementsList = []
    for i in range(len(lst)):
        Elements = lst[i].getElementsByTagName(type_of_element)
        for elm in Elements:
            elementsList.append(elm)
    return elementsList


def ParaToText(p):
    rs = p._element.xpath('.//w:t')
    return u" ".join([r.text for r in rs])


# ------RÚT TRÍCH TEXT TỪ FILE----#
# Docx
def DocxToText(docxFileName):
    # Parse xml file
    xmlFileName = 'mydocx.xml'
    opc_to_flat_opc(docxFileName, xmlFileName)
    myDocx = minidom.parse(xmlFileName)

    # Get elements
    paragraph = myDocx.getElementsByTagName('w:p')
    table = myDocx.getElementsByTagName('w:tbl')

    # Get all w:p elements in table elements. Output is two-dimensional list
    wpTbl = GetAllElements(table, 'w:p')

    # Get text and save to "string" variable
    paraIndex = 0
    tblIndex = 0
    string = ""
    while paraIndex < len(paragraph):
        if paragraph[paraIndex] in wpTbl:
            string = string + TableString(table[tblIndex])
            paraIndex += len(table[tblIndex].getElementsByTagName('w:p'))
            tblIndex += 1

        else:
            string = string + ParaString(paragraph[paraIndex]) + '\n'
            paraIndex += 1

        string = string + '\n'
    os.remove("mydocx.xml")
    return string


# Doc
def DocToDocx(filename, path=os.getcwd()):
    baseDir = os.path.abspath(os.getcwd())
    pythoncom.CoInitialize()
    word = win32com.client.Dispatch("Word.application")
    filePath = os.path.join(baseDir, filename)
    fileName, fileExtension = os.path.splitext(filePath)

    if "~$" not in fileName:
        if fileExtension.lower() == '.doc':
            extension = str(uuid.uuid4().hex[:10]).format(filePath, 'x')
            docxFile = fileName + extension
            if not os.path.isfile(docxFile):

                filePath = os.path.abspath(filePath)
                docxFile = os.path.abspath(docxFile)

                try:
                    wordDoc = word.Documents.Open(filePath)
                    wordDoc.SaveAs2(docxFile, FileFormat=16)
                    wordDoc.Close()
                except Exception as e:
                    print('Failed to Convert: {0}'.format(filePath))
                    print(e)
            return docxFile + ".docx"


# Powerpoint
def PptsToText(filename):
    ppt = Presentation(filename)
    string = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                string += shape.text + "\n"
    return string


# CSV
def CsvToText(filename):
    string = ""
    with open(filename, 'rt', encoding="utf-8") as f:
        data = csv.reader(f)
        for row in data:
            for cell in row:
                if cell != '':
                    string += cell + '\n'
    return string


# Excel
# hàm rút trich các câu từ file excel ### cần cập nhật thêm vì chưa hoàn thiện
def XlsxToText(filename):
    fileName, fileExtension = os.path.splitext(filename)
    data = pd.read_excel(filename, index_col=0, keep_default_na=0)
    filename_csv = fileName + str(uuid.uuid4().hex[:10]) + ".csv"
    data.to_csv(filename_csv, encoding='utf-8')
    res = CsvToText(filename_csv)
    os.remove(filename_csv)
    return res


# Tách đoạn
def ListPara(document):
    return document.split('\n')


# Tách câu
def ParaToSen(para):
    return sent_tokenize(para)


def ListSentence(para_list):
    sentences = []

    for p in para_list:
        sentences.extend(ParaToSen(p))

    return sentences


# Tách từ
# Giữa các chữ trong 1 từ CÓ dấu gạch dưới (underscore)
def SentenceToWordUnderscore(sentence):
    words = word_tokenize(sentence)

    for i in range(len(words)):
        words[i] = words[i].replace(' ', '_')

    return words


# Giữa các chữ trong 1 từ KHÔNG CÓ dấu gạch dưới (underscore)
def SentenceToWord(sentence):
    return word_tokenize(sentence)


def ListWord(para_list):
    sentences = ListSentence(para_list)

    words = []
    for s in sentences:
        words.append(SentenceToWord(s))

    return words


# Số từ của mỗi câu trong văn bản
def NumOfWord(words_list):
    lst = []

    for sent in words_list:
        lst.append(len(sent))

    return lst


# Ghép các từ thành 1 câu
# Input: word_list [word1, word2, word3,...]
# Output: string
def JoinWord(word_list):
    string = ""

    for i in range(len(word_list) - 1):
        string += word_list[i]
        if word_list[i + 1] not in punctuation:
            string += ' '
    string += word_list[len(word_list) - 1]

    return string


# Chia nhỏ câu nếu câu có độ dài > max_len (nhập thủ công)
# Input:
#   + sentence (string)
#   + max_len: độ dài tối đa của 1 câu
# Output: list các sentence mới (string)
def SplitSen(sentence, max_len):
    lst = []

    words = SentenceToWordUnderscore(sentence)
    newSen = JoinWord(words)

    x = math.ceil(len(newSen) / max_len)
    newLen = math.ceil(len(newSen) / x)

    start = 0
    end = newLen
    for i in range(x - 1):
        string = newSen[start: end]

        index = end - 1
        if newSen[index] != ' ':
            for j in range(index + 1, len(newSen)):
                if newSen[j] == ' ':
                    break
                string += newSen[j]
            start = j + 1
            end = start + newLen

        else:
            start += newLen
            end = start + newLen

        lst.append(string.strip().replace('_', ' '))
        string = ""

    lst.append(newSen[start:].strip().replace('_', ' '))

    return lst


# Input: list sentences
# Output: list các sentence mới thỏa điều kiện chia nhỏ và sentence cũ
def SplitSenList(sent_list, max_len):
    newLst = []
    minLen = 5  # cas cau co it hon 5 tu
    for sent in sent_list:
        if len(sent) > max_len:
            newLst.extend(SplitSen(sent, max_len))
        elif len(sent) > minLen:
            newLst.append(sent)

    return newLst


# ---------TIỀN XỬ LÝ--------#


def Preprocess(filename):
    name, fileExtension = os.path.splitext(filename)
    para = []

    extension = [".doc", ".docx", ".pdf", ".xlsx", ".csv", ".pptx", ".txt"]

    if fileExtension.lower() not in extension:
        raise TypeError("Wrong type document file!")

    else:
        if fileExtension.lower() == ".doc":
            newFilenameDocx = DocToDocx(filename)
            doc = DocxToText(newFilenameDocx)

        if fileExtension.lower() == ".docx":
            doc = DocxToText(filename)

        if fileExtension.lower() == ".pdf":
            para = Pdf_extract.pdf2txt(filename)

        if fileExtension.lower() == ".xlsx":
            doc = XlsxToText(filename)

        if fileExtension.lower() == ".csv":
            doc = CsvToText(filename)

        if fileExtension.lower() == ".pptx":
            doc = PptsToText(filename)

        if fileExtension.lower() == ".txt":
            f = open(filename, 'r', encoding='utf-8')
            doc = f.read()
            f.close()

    if fileExtension.lower() != ".pdf":
        para = ListPara(doc)

    listSenTemp = SplitSenList(ListSentence(para), 450)
    listSentence = [sen.replace("\xa0", "") for sen in listSenTemp]
    ListWords = ListWord(para)

    return listSentence, ListWords, os.path.basename(filename)


if __name__ == '__main__':
    FILE_PATH = 'project_doc.docx'

    start_time = time.time()

    sent, word, filename = Preprocess(FILE_PATH)
    print("ds câu: ", sent, "\n\n\n word: ", word, "\n\n\nfilename", filename)
    end_time = time.time()

    print("Time:", end_time - start_time)
