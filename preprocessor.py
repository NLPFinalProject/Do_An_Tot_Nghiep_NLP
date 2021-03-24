import tika
from tika import parser
import re
import os
import sys
from pptx import Presentation
import pandas as pd
import openpyxl
from vncorenlp import VnCoreNLP
import zipfile
import xml.etree.ElementTree
import logging
from docx_utils.flatten import opc_to_flat_opc
from xml.dom import minidom
import win32com.client
import docx
import PyPDF2
import uuid
from os.path import dirname, abspath, join

# LƯU Ý:
# input của các hàm xử lý là tên file và output của các hàm xử lý ...2txt là một list các câu.
# từng phần tử của 1 list lớn (1 câu) là 1 list nhỏ chứa các từ được phân tách nhờ VNCORE.

#### ------------------------------------các function hỗ trợ cho function docx2txt
# Get paragraph string. Input is paragraph element
def para_string(para):
    string = ""
    # if (str(para)[21:34] not in str(wp_tbl)):# and (str(para)[21:34] not in str(wp_txbx)):
    wt = para.getElementsByTagName('w:t')
    for i in range(len(wt)):
        string = string + wt[i].firstChild.data

    return string


# Get table string. Input is table element
def table_string(table):
    string = ""

    wp = table.getElementsByTagName('w:p')
    column = len(table.getElementsByTagName('w:tc')) / len(table.getElementsByTagName('w:tr'))
    c = 1
    for i in range(len(wp)):
        string = string + para_string(wp[i])
        if c % column == 0 and c != len(wp):
            string = string + '\n'
        else:
            string = string + '\t'
        c = c + 1
    return string


# Get all elements
def get_all_elements(lst, type_of_element):
    elements_list = []
    for i in range(len(lst)):
        Elements = lst[i].getElementsByTagName(type_of_element)
        for elm in Elements:
            elements_list.append(elm)

    return elements_list


def para2text(p):
    rs = p._element.xpath('.//w:t')
    return u" ".join([r.text for r in rs])


##--------------------------------------------------------------------

def docx2txt(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para2text(para))
    for table in doc.tables:
        for row in table.rows:
            s = ""
            for cell in row.cells:
                if (cell.text == ""):
                    continue
                s += cell.text + '. '
            fullText.append(s)
    # return '\n'.join(fullText),
    if("preprocessor" in os.getcwd()):
        vncorenlp_file = os.getcwd()+'/VnCoreNLP/VnCoreNLP-1.1.1.jar'
    else:
        vncorenlp_file = os.getcwd()+'/preprocessor/VnCoreNLP/VnCoreNLP-1.1.1.jar'
    split_sentence = []
    with VnCoreNLP(vncorenlp_file, annotators="wseg", max_heap_size='-Xmx4g', quiet=False) as vncorenlp:
        for sentences in fullText:
            split_sentence.extend(vncorenlp.tokenize(sentences))
    return split_sentence


# hàm chuyển đổi định dạng từ fild .doc sang .docx
def doc2docx(filename, path=os.getcwd()):
    baseDir = os.path.abspath(os.getcwd())  # Starting directory for directory walk
    word = win32com.client.Dispatch("Word.application")
    file_path = os.path.join(baseDir, filename)
    file_name, file_extension = os.path.splitext(file_path)

    if "~$" not in file_name:
        if file_extension.lower() == '.doc':  #
            # docx_file = '{0}{1}'.format(file_path, 'x')
            docx_file = file_name + str(uuid.uuid4().hex[:10]).format(file_path, 'x')
            if not os.path.isfile(docx_file):  # Skip conversion where docx file already exists

                file_path = os.path.abspath(file_path)
                docx_file = os.path.abspath(docx_file)

                try:
                    wordDoc = word.Documents.Open(file_path)
                    wordDoc.SaveAs2(docx_file, FileFormat=16)
                    wordDoc.Close()
                except Exception as e:
                    print('Failed to Convert: {0}'.format(file_path))
                    print(e)
            return docx_file+".docx"


# hàm rút trích các câu từ file ppt
def ppt2txt(filename):
    ppt = Presentation(filename)
    sentences = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                sentences += shape.text + ". "

    if("preprocessing" in os.getcwd()):
        vncorenlp_file = os.getcwd()+'/VnCoreNLP/VnCoreNLP-1.1.1.jar'
    else:
        vncorenlp_file = os.getcwd()+'/preprocessing/VnCoreNLP/VnCoreNLP-1.1.1.jar'
    with VnCoreNLP(vncorenlp_file, annotators="wseg", max_heap_size='-Xmx4g', quiet=False) as vncorenlp:
        split_sentence = vncorenlp.tokenize(sentences)
    return split_sentence


# hàm rút trich các câu từ file excel ### cần cập nhật thêm vì chưa hoàn thiện
def xlsx2txt(filename):
    data = pd.read_excel(filename, sheet_name=None, index_col=0, keep_default_na=0)
    list_sheet = []
    for key, value in data.items():
        list_sheet.append(value)
    listToStr = '.'.join(map(str, list_sheet))
    list_sen = listToStr.split('\n')
    for i in range(0, len(list_sen)):
        list_sen[i] = " ".join(list_sen[i].split())
    return list_sen


def convert2listsen(vncoretoken):
    # remove .
    for i in vncoretoken:
        for j in i:
            if (j == "."):
                i.remove(j)
    s = ""
    res = []
    for i in vncoretoken:
        for j in i:
            s += j + " "
        s = s.strip()
        s += "."
        res.append(s)
        s = ""
    return res


def num_of_word(list_sentences):
    num_word = []
    for i in list_sentences:
        num_word.append(len(i))
    return num_word


def preprocess(filename):
    # thực hiện đưa filename nào muốn xử lý vào biến filename bên dưới và đợi kết quả trên màn hình.
    name, file_extension = os.path.splitext(filename)
    if (file_extension.lower() == ".doc"):
        new_filename_docx = doc2docx(filename)
        a = docx2txt(new_filename_docx)
        os.remove(new_filename_docx)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".docx"):
        a = docx2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)  # số từ của câu đầu tiên tương tự cho a[1],....

    # elif (file_extension.lower() == ".pdf"):
    #     a = pdf2txt(filename)
    #     b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
    #     num_word = num_of_word(b)

    elif (file_extension.lower() == ".xlsx"):
        a = xlsx2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)

    elif (file_extension.lower() == ".pptx"):
        a = ppt2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)

    return filename, b, num_word  # filename, list câu. số từ của mỗi câu





if __name__ == '__main__':
    filename = "sample.doc"
    a, b, c = preprocess(filename)
    print("Tên file là: ",a)
    print("\n Danh sách các câu của file là: ",b)
    print("\n Danh sách số từ của file là: ",c)
    
